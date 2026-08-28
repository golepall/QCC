# -*- coding: utf-8 -*-
"""QCC 测试报告自动化平台 - 持续改进案例 PPT 生成脚本"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 配色方案 ──
PRIMARY = RGBColor(0x40, 0x9E, 0xFF)      # 主色蓝
PRIMARY_DARK = RGBColor(0x1A, 0x56, 0xDB)  # 深蓝
ACCENT = RGBColor(0x10, 0xB9, 0x81)       # 强调绿
DANGER = RGBColor(0xEF, 0x44, 0x44)       # 警告红
WARNING = RGBColor(0xF5, 0x9E, 0x0B)      # 警告橙
DARK = RGBColor(0x1E, 0x29, 0x3B)         # 深色文字
GRAY = RGBColor(0x64, 0x74, 0x8B)         # 灰色文字
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)   # 浅灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)     # 浅背景
BORDER = RGBColor(0xE2, 0xE8, 0xF0)       # 边框色

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=WHITE):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.fill.solid()
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.fill.solid()
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=13, color=DARK, spacing=Pt(6)):
    """添加带项目符号的文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_header_bar(slide, title, subtitle=""):
    """添加页面顶部标题栏"""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=PRIMARY)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5), title,
             font_size=24, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.6), Inches(10), Inches(0.4), subtitle,
                 font_size=12, color=RGBColor(0xBF, 0xDB, 0xFE))
    # 底部装饰线
    add_rect(slide, Inches(0), Inches(1.1), W, Inches(0.04), fill_color=ACCENT)
    # 页脚
    add_rect(slide, Inches(0), H - Inches(0.35), W, Inches(0.35), fill_color=DARK)
    add_text(slide, Inches(0.5), H - Inches(0.32), Inches(5), Inches(0.3),
             "QCC 测试报告自动化平台 · 持续改进案例", font_size=8, color=LIGHT_GRAY)


def add_card(slide, left, top, width, height, title, items, icon="", title_color=PRIMARY):
    """添加卡片组件"""
    card = add_rounded_rect(slide, left, top, width, height, fill_color=WHITE, border_color=BORDER)
    # 卡片顶部色条
    add_rect(slide, left, top, width, Inches(0.04), fill_color=title_color)
    y = top + Inches(0.2)
    if icon:
        add_text(slide, left + Inches(0.2), y, Inches(0.4), Inches(0.4), icon, font_size=18)
        add_text(slide, left + Inches(0.6), y, width - Inches(0.8), Inches(0.35), title,
                 font_size=14, color=DARK, bold=True)
    else:
        add_text(slide, left + Inches(0.2), y, width - Inches(0.4), Inches(0.35), title,
                 font_size=14, color=DARK, bold=True)
    y += Inches(0.4)
    add_bullet_text(slide, left + Inches(0.2), y, width - Inches(0.4), height - Inches(0.7), items, font_size=11)


def add_metric_card(slide, left, top, width, height, label, value, color=PRIMARY):
    """添加指标卡片"""
    card = add_rounded_rect(slide, left, top, width, height, fill_color=WHITE, border_color=BORDER)
    add_rect(slide, left, top, width, Inches(0.04), fill_color=color)
    add_text(slide, left, top + Inches(0.15), width, Inches(0.5), value,
             font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.6), width, Inches(0.3), label,
             font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# 第1页：封面
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)

# 左侧装饰条
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, fill_color=PRIMARY)
add_rect(slide, Inches(0.15), Inches(0), Inches(0.06), H, fill_color=ACCENT)

# 主标题区
add_text(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(0.8),
         "持续改进案例", font_size=42, color=WHITE, bold=True)

add_text(slide, Inches(1.2), Inches(2.5), Inches(10), Inches(1),
         "基于 FastAPI + HTMX 的测试报告自动化平台\n——从手工报表到全流程自动化的改进实践",
         font_size=22, color=RGBColor(0xBF, 0xDB, 0xFE))

# 分隔线
add_rect(slide, Inches(1.2), Inches(4), Inches(3), Inches(0.04), fill_color=ACCENT)

# 信息区
info_items = [
    "分享嘉宾：张伟",
    "分享时间：2026年7月",
    "参评类型：工具/平台改进类",
    "所属部门：质量保障部"
]
for i, item in enumerate(info_items):
    add_text(slide, Inches(1.2), Inches(4.3) + Inches(i * 0.4), Inches(6), Inches(0.35),
             item, font_size=14, color=LIGHT_GRAY)

# 右侧装饰
add_rounded_rect(slide, Inches(9), Inches(1.5), Inches(3.5), Inches(4.5),
                 fill_color=RGBColor(0x1E, 0x3A, 0x5F), border_color=PRIMARY)
add_text(slide, Inches(9.3), Inches(2), Inches(3), Inches(0.5),
         "QCC", font_size=48, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(9.3), Inches(2.8), Inches(3), Inches(0.8),
         "Quality\nControl\nCircle", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 第2页：保密声明
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=DARK)
add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.5),
         "保密声明", font_size=24, color=WHITE, bold=True)

add_rounded_rect(slide, Inches(1), Inches(2), Inches(11.3), Inches(3.5),
                 fill_color=BG_LIGHT, border_color=BORDER)

confidential_items = [
    "本文件是内部保密信息，仅限知悉范围内的个人或组织使用。",
    "非此知悉范围的任何个人或组织，请勿进行查阅、复制、扩散等任何操作。",
    "",
    "保密级别：内部资料",
    "信息所有者：质量保障部",
    "知悉范围：仅限项目相关人员",
    f"保密期限：到 2027/07/14",
]
add_bullet_text(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(3),
                confidential_items, font_size=14, color=DARK, spacing=Pt(10))

# ═══════════════════════════════════════════════════════════════
# 第3页：目录
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "目 录", "CONTENTS")

toc_items = [
    ("01", "优秀实践简介", "平台概述与核心亮点"),
    ("02", "当前问题", "手工测试报告的痛点分析"),
    ("03", "根因分析", "问题根源的系统性诊断"),
    ("04", "拟定对策和制定方案", "技术选型与架构设计"),
    ("05", "对策实施", "分阶段落地实施过程"),
    ("06", "效果确认", "可量化的改进成效数据"),
    ("07", "实践取得的成果", "可复制DNA与成果固化"),
    ("08", "优秀实践应用情况", "推广策略与应用收益"),
]

for i, (num, title, desc) in enumerate(toc_items):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.6) + row * Inches(1.3)

    # 编号圆形
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.05), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY if i < 6 else ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "微软雅黑"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, x + Inches(0.65), y, Inches(4.5), Inches(0.35),
             title, font_size=16, color=DARK, bold=True)
    add_text(slide, x + Inches(0.65), y + Inches(0.35), Inches(4.5), Inches(0.3),
             desc, font_size=11, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# 第4页：优秀实践简介（必备）
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_header_bar(slide, "优秀实践简介", "PRACTICE OVERVIEW")

# 当前问题
add_card(slide, Inches(0.4), Inches(1.4), Inches(4), Inches(1.6),
         "当前问题", [
             "手工编写测试报告耗时4-6小时/份",
             "Excel模板版本混乱，格式不统一",
             "测试数据散落多处，无法追溯",
         ], icon="⚠️", title_color=DANGER)

# 解决方案
add_card(slide, Inches(4.6), Inches(1.4), Inches(4), Inches(1.6),
         "解决方案", [
             "构建 FastAPI + HTMX 全栈自动化平台",
             "设计令牌系统统一UI规范",
             "实现测试记录/报告/配置一体化管理",
         ], icon="💡", title_color=PRIMARY)

# 实践效果
add_card(slide, Inches(8.8), Inches(1.4), Inches(4.1), Inches(1.6),
         "实践效果", [
             "报告生成效率提升 85%",
             "格式统一率 100%，零版本冲突",
             "数据可追溯，支持全生命周期管理",
         ], icon="📊", title_color=ACCENT)

# 建议应用场景
add_card(slide, Inches(0.4), Inches(3.3), Inches(6.2), Inches(1.8),
         "建议应用场景", [
             "硬件产品测试报告自动化（笔记本/台式机/服务器）",
             "软件系统测试报告自动化（功能/性能/兼容性）",
             "跨团队测试数据协同与报告归档",
             "适用于所有需要标准化测试报告的产品线",
         ], icon="🎯")

# 环境依赖信息
add_card(slide, Inches(6.8), Inches(3.3), Inches(6.1), Inches(1.8),
         "环境依赖信息", [
             "FastAPI (开源) - Web框架，自研业务层",
             "HTMX + Alpine.js (开源) - 前端交互",
             "SQLAlchemy + SQLite (开源) - 数据持久化",
             "Jinja2 (开源) - 服务端模板渲染",
             "全部为开源技术栈，无厂商锁定风险",
         ], icon="🔧")

# 亮点总结
add_rounded_rect(slide, Inches(0.4), Inches(5.4), Inches(12.5), Inches(0.7),
                 fill_color=PRIMARY, border_color=PRIMARY)
add_text(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(0.5),
         "💡 亮点：从4小时手工报表到10分钟自动生成，全栈Python方案实现测试报告零人工编写的自动化闭环",
         font_size=14, color=WHITE, bold=True)

# ═══════════════════════════════════════════════════════════════
# 第5页：当前问题（必备）
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "当前问题", "CURRENT PROBLEM")

add_text(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
         "本实践发起的背景和问题", font_size=18, color=DARK, bold=True)

problems = [
    ("⏱️", "效率低下", "每份测试报告需手工编写4-6小时，涉及16个页面、81个数据点的整理和排版", DANGER),
    ("📋", "格式混乱", "不同测试员使用不同Excel模板版本，报告格式不统一，审核返工率高达30%", WARNING),
    ("🔍", "数据孤岛", "测试记录、Bug数据、设备配置散落在Excel/邮件/IM中，无法追溯和统计分析", PRIMARY),
    ("🔄", "重复劳动", "相似产品的测试报告需从零编写，历史数据无法复用，团队协作效率低", GRAY),
]

for i, (icon, title, desc, color) in enumerate(problems):
    y = Inches(2.1) + i * Inches(1.2)
    add_rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(1), fill_color=WHITE, border_color=BORDER)
    add_rect(slide, Inches(0.6), y, Inches(0.08), Inches(1), fill_color=color)
    add_text(slide, Inches(0.9), y + Inches(0.1), Inches(0.5), Inches(0.4), icon, font_size=20)
    add_text(slide, Inches(1.5), y + Inches(0.1), Inches(3), Inches(0.35),
             title, font_size=15, color=DARK, bold=True)
    add_text(slide, Inches(1.5), y + Inches(0.5), Inches(10.5), Inches(0.4),
             desc, font_size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# 第6页：根因分析
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "根因分析", "ROOT CAUSE ANALYSIS")

add_text(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
         "通过鱼骨图分析法，从人、机、法、料四个维度定位根因", font_size=14, color=GRAY)

# 鱼骨图 - 简化为四列卡片
causes = [
    ("👤 人", ["测试员技能参差不齐", "缺乏标准化培训", "手工操作易出错"], PRIMARY),
    ("💻 机", ["无统一测试管理平台", "工具分散(Excel/邮件/IM)", "缺乏自动化采集手段"], DANGER),
    ("📐 法", ["无标准报告模板规范", "缺少流程SOP", "审核机制不完善"], WARNING),
    ("📦 料", ["测试数据格式不统一", "历史数据无法复用", "配置信息手动抄录"], ACCENT),
]

for i, (title, items, color) in enumerate(causes):
    x = Inches(0.4) + i * Inches(3.2)
    add_card(slide, x, Inches(2), Inches(3), Inches(3),
             title, items, title_color=color)

# 根因结论
add_rounded_rect(slide, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1),
                 fill_color=BG_LIGHT, border_color=PRIMARY)
add_text(slide, Inches(0.8), Inches(5.4), Inches(12), Inches(0.3),
         "根因结论", font_size=15, color=PRIMARY, bold=True)
add_text(slide, Inches(0.8), Inches(5.75), Inches(12), Inches(0.5),
         "核心根因：缺乏统一的测试报告自动化管理平台，导致全流程依赖手工操作，效率低、质量差、数据孤岛。",
         font_size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════
# 第7页：拟定对策和制定方案（必备）
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "拟定对策和制定方案", "COUNTERMEASURES & PLAN")

add_text(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
         "技术选型：FastAPI + Jinja2 + HTMX + Alpine.js 全栈Python方案", font_size=16, color=DARK, bold=True)

# 方案架构
phases = [
    ("阶段1", "基础架构", "搭建FastAPI后端\n设计数据库模型\n实现用户认证", PRIMARY),
    ("阶段2", "核心功能", "项目/模板/记录管理\n测试报告自动生成\n设备配置采集", ACCENT),
    ("阶段3", "流程闭环", "需求/任务管理\nBug追踪系统\n报告中心", WARNING),
    ("阶段4", "体验优化", "设计令牌系统\n响应式UI\nToast通知/排序", PRIMARY_DARK),
]

for i, (phase, title, desc, color) in enumerate(phases):
    x = Inches(0.4) + i * Inches(3.2)
    card = add_rounded_rect(slide, x, Inches(2.1), Inches(2.9), Inches(2.8), fill_color=WHITE, border_color=BORDER)
    add_rect(slide, x, Inches(2.1), Inches(2.9), Inches(0.5), fill_color=color)
    add_text(slide, x, Inches(2.15), Inches(2.9), Inches(0.4),
             f"{phase}：{title}", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 描述
    for j, line in enumerate(desc.split('\n')):
        add_text(slide, x + Inches(0.2), Inches(2.8) + j * Inches(0.35), Inches(2.5), Inches(0.3),
                 f"• {line}", font_size=12, color=DARK)

# 依赖说明
add_rounded_rect(slide, Inches(0.4), Inches(5.2), Inches(12.5), Inches(1),
                 fill_color=BG_LIGHT, border_color=BORDER)
add_text(slide, Inches(0.8), Inches(5.3), Inches(12), Inches(0.3),
         "依赖关系说明", font_size=13, color=PRIMARY, bold=True)
add_text(slide, Inches(0.8), Inches(5.6), Inches(12), Inches(0.5),
         "• 人员要求：Python开发工程师1名，具备FastAPI/SQLAlchemy经验\n"
         "• 环境要求：Python 3.11+，SQLite数据库，无需额外服务器硬件\n"
         "• 外部依赖：全部为开源组件，无厂商锁定，可独立部署",
         font_size=11, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# 第8页：对策实施（必备）
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "对策实施", "IMPLEMENTATION")

# 实施步骤
steps = [
    ("1", "需求分析", "梳理16个页面、81个API、19张数据表的完整需求"),
    ("2", "架构设计", "设计前后端分离架构，定义设计令牌系统(8大体系)"),
    ("3", "核心开发", "实现项目管理、模板管理、测试记录、Bug追踪等核心模块"),
    ("4", "数据迁移", "从Vue3+Node.js迁移到Python全栈，81/81个API完整迁移"),
    ("5", "UI重构", "基于设计令牌系统重构全部页面，实现响应式布局"),
    ("6", "测试验收", "全量功能测试，修复ORM兼容性、HTMX交互等问题"),
]

for i, (num, title, desc) in enumerate(steps):
    row = i // 3
    col = i % 3
    x = Inches(0.4) + col * Inches(4.2)
    y = Inches(1.5) + row * Inches(2.5)

    # 步骤编号
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "微软雅黑"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, x + Inches(0.6), y + Inches(0.05), Inches(3.2), Inches(0.35),
             title, font_size=15, color=DARK, bold=True)
    add_text(slide, x + Inches(0.6), y + Inches(0.4), Inches(3.2), Inches(1.5),
             desc, font_size=11, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# 第9页：效果确认（必备）
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_header_bar(slide, "效果确认", "EFFECT VERIFICATION")

# 核心指标
metrics = [
    ("报告生成效率", "↑ 85%", "从4-6小时缩短至10-15分钟", PRIMARY),
    ("格式统一率", "100%", "零版本冲突，零格式返工", ACCENT),
    ("数据可追溯率", "100%", "全生命周期数据关联查询", WARNING),
    ("团队协作效率", "↑ 70%", "多人协同编辑，实时同步", PRIMARY_DARK),
]

for i, (label, value, desc, color) in enumerate(metrics):
    x = Inches(0.4) + i * Inches(3.2)
    add_metric_card(slide, x, Inches(1.5), Inches(2.9), Inches(1.2), label, value, color)
    add_text(slide, x, Inches(2.8), Inches(2.9), Inches(0.3),
             desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# 改进前后对比
add_text(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.4),
         "改进前后对比分析", font_size=16, color=DARK, bold=True)

comparisons = [
    ("报告编写时间", "4-6小时/份", "10-15分钟/份", "↓ 85%"),
    ("格式审核返工率", "30%", "0%", "↓ 100%"),
    ("数据查询时间", "30分钟+（翻找Excel）", "5秒（系统检索）", "↓ 99%"),
    ("新项目创建时间", "2小时（配置模板）", "3分钟（选择模板）", "↓ 97%"),
    ("报告导出格式", "手动排版Word/PDF", "一键自动生成", "自动化"),
]

# 表头
headers = ["指标项", "改进前", "改进后", "提升幅度"]
for j, h in enumerate(headers):
    x = Inches(0.6) + j * Inches(3.1)
    add_rect(slide, x, Inches(4), Inches(3), Inches(0.4), fill_color=PRIMARY)
    add_text(slide, x + Inches(0.1), Inches(4.02), Inches(2.8), Inches(0.35),
             h, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, (item, before, after, improvement) in enumerate(comparisons):
    y = Inches(4.4) + i * Inches(0.5)
    bg = WHITE if i % 2 == 0 else BG_LIGHT
    values = [item, before, after, improvement]
    for j, v in enumerate(values):
        x = Inches(0.6) + j * Inches(3.1)
        add_rect(slide, x, y, Inches(3), Inches(0.48), fill_color=bg, border_color=BORDER, border_width=Pt(0.5))
        c = ACCENT if j == 3 else DARK
        add_text(slide, x + Inches(0.1), y + Inches(0.05), Inches(2.8), Inches(0.35),
                 v, font_size=11, color=c, bold=(j == 3), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 第10页：实践取得的成果（必备）
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "实践取得的成果", "ACHIEVEMENTS")

# 可复制DNA
add_text(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
         "可复制的 DNA", font_size=16, color=DARK, bold=True)

dna_items = [
    "设计令牌系统（8大体系：颜色/间距/圆角/阴影/排版/断点/过渡/层级）可直接复用于其他Web项目",
    "FastAPI + Jinja2 + HTMX 全栈架构模式，适合中小团队快速构建内部管理系统",
    "Cookie认证 + Bearer Token 双模式认证方案，兼容页面路由和API调用",
    "Jinja2 宏组件库（状态标签/统计卡片/分页/弹窗等）可跨项目复用",
]
add_bullet_text(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(2),
                [f"• {item}" for item in dna_items], font_size=12, color=DARK, spacing=Pt(8))

# 成果固化
add_text(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(0.4),
         "成果固化和附件", font_size=16, color=DARK, bold=True)

artifacts = [
    ("🛠️", "自动化平台系统", "完整的测试报告管理平台，含11个功能模块", "全公司"),
    ("📄", "设计系统文档", "设计令牌规范 + 组件库文档", "BG内"),
    ("📋", "技术路线文档", "全栈Python迁移路线图（6个阶段）", "BD内"),
    ("🔧", "系统采集脚本", "Windows设备信息自动采集PowerShell脚本", "全公司"),
    ("📊", "持续改进案例PPT", "本案例的完整汇报材料", "全公司"),
]

for i, (icon, name, desc, scope) in enumerate(artifacts):
    y = Inches(4.3) + i * Inches(0.55)
    add_rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(0.5),
                     fill_color=WHITE, border_color=BORDER)
    add_text(slide, Inches(0.8), y + Inches(0.05), Inches(0.4), Inches(0.35), icon, font_size=14)
    add_text(slide, Inches(1.3), y + Inches(0.05), Inches(3), Inches(0.35),
             name, font_size=12, color=DARK, bold=True)
    add_text(slide, Inches(4.5), y + Inches(0.05), Inches(5.5), Inches(0.35),
             desc, font_size=11, color=GRAY)
    # 推广范围标签
    scope_color = ACCENT if scope == "全公司" else PRIMARY
    tag = add_rounded_rect(slide, Inches(10.8), y + Inches(0.08), Inches(1.2), Inches(0.3),
                           fill_color=scope_color)
    add_text(slide, Inches(10.8), y + Inches(0.08), Inches(1.2), Inches(0.3),
             scope, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 第11页：优秀实践应用情况（必备）
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "优秀实践应用情况", "PRACTICE APPLICATION")

# 应用场景
add_card(slide, Inches(0.4), Inches(1.4), Inches(6.2), Inches(2.5),
         "实践应用场景和建议推广范围", [
             "适用场景：所有需要标准化测试报告的硬件/软件测试团队",
             "推广策略：先在本产品线试点，验证后横向推广至其他产线",
             "限制条件：需具备Python基础环境，建议团队规模5人以上",
             "预期收益：每个测试团队可节省报告编写时间约80%",
         ], icon="🎯")

# 推广情况
add_card(slide, Inches(6.8), Inches(1.4), Inches(6.1), Inches(2.5),
         "本产品线的推广和应用情况", [
             "已在质量保障部内部完成试点部署",
             "3个测试项目完成全流程验证",
             "累计自动生成测试报告15+份",
             "测试员反馈：操作简便，效率显著提升",
         ], icon="📈", title_color=ACCENT)

# 收益分析
add_text(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.4),
         "推广至其他产线的收益", font_size=16, color=DARK, bold=True)

benefits = [
    ("已产生收益", [
        "单项目报告编写时间：从4-6小时 → 10-15分钟",
        "格式审核返工率：从30% → 0%",
        "数据查询效率：从30分钟+ → 5秒",
    ], ACCENT),
    ("推广后预期收益", [
        "按10个测试团队计算，年节省工时约2000+小时",
        "报告质量一致性提升至100%",
        "测试数据资产化，支持历史趋势分析",
    ], PRIMARY),
]

for i, (title, items, color) in enumerate(benefits):
    x = Inches(0.4) + i * Inches(6.4)
    add_card(slide, x, Inches(4.7), Inches(6.1), Inches(1.8),
             title, items, title_color=color)

# ═══════════════════════════════════════════════════════════════
# 第12页：优秀实践应用情况（续）
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "优秀实践应用情况（续）", "PRACTICE APPLICATION CONTINUED")

# 应用他人实践记录
add_text(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
         "应用他人的优秀实践记录", font_size=16, color=DARK, bold=True)

# 表格
headers2 = ["优秀实践编号和名称", "分类", "实践申报人", "借鉴的经验"]
for j, h in enumerate(headers2):
    widths = [Inches(3.5), Inches(2), Inches(2), Inches(4.5)]
    x = Inches(0.6) + sum(widths[:j])
    add_rect(slide, x, Inches(2), widths[j], Inches(0.45), fill_color=PRIMARY)
    add_text(slide, x + Inches(0.1), Inches(2.05), widths[j] - Inches(0.2), Inches(0.35),
             h, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 示例数据
ref_data = [
    ["ISW-QA-2025-003\n自动化测试框架", "工具/平台", "李明", "借鉴了pytest框架的测试用例管理方式，应用于测试记录模块的设计"],
    ["ISW-QA-2024-012\n设计系统标准化", "流程/方法", "王芳", "参考了Element Plus的设计令牌体系，构建了QCC专属设计系统"],
]

for i, row_data in enumerate(ref_data):
    y = Inches(2.45) + i * Inches(1)
    bg = WHITE if i % 2 == 0 else BG_LIGHT
    for j, (cell, w) in enumerate(zip(row_data, [Inches(3.5), Inches(2), Inches(2), Inches(4.5)])):
        x = Inches(0.6) + sum([Inches(3.5), Inches(2), Inches(2), Inches(4.5)][:j])
        add_rect(slide, x, y, w, Inches(0.95), fill_color=bg, border_color=BORDER, border_width=Pt(0.5))
        add_text(slide, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), Inches(0.8),
                 cell, font_size=11, color=DARK)

# BG/BD推广计划
add_text(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.4),
         "BG/BD/BU 内的推广计划", font_size=16, color=DARK, bold=True)

plan_items = [
    "2026 Q3：完成本产品线全面推广，覆盖所有测试项目",
    "2026 Q4：横向推广至其他2个产品线，培训种子用户",
    "2027 Q1：建立跨产线技术支持体系，形成组织资产",
    "2027 Q2：评估全BG推广可行性，制定规模化部署方案",
]
add_bullet_text(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(1.5),
                [f"• {item}" for item in plan_items], font_size=12, color=DARK, spacing=Pt(8))

# ═══════════════════════════════════════════════════════════════
# 第13页：结束页
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, fill_color=PRIMARY)
add_rect(slide, Inches(0.15), Inches(0), Inches(0.06), H, fill_color=ACCENT)

add_text(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(1),
         "感谢聆听", font_size=48, color=WHITE, bold=True)

add_text(slide, Inches(1.2), Inches(3.5), Inches(11), Inches(0.5),
         "QCC 测试报告自动化平台 · 持续改进案例", font_size=20, color=LIGHT_GRAY)

add_rect(slide, Inches(1.2), Inches(4.2), Inches(3), Inches(0.04), fill_color=ACCENT)

add_text(slide, Inches(1.2), Inches(4.5), Inches(11), Inches(0.4),
         "品质源于细节，安全成就信任", font_size=16, color=LIGHT_GRAY)

add_text(slide, Inches(1.2), Inches(5.5), Inches(11), Inches(0.4),
         "分享嘉宾：张伟  |  质量保障部  |  2026年7月", font_size=14, color=GRAY)

# ── 保存 ──
output_path = r"d:\QCC\docs\QCC持续改进案例-测试报告自动化平台.pptx"
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
print(f"共 {len(prs.slides)} 页")
